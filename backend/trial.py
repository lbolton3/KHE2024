from pptx import Presentation
import google.generativeai as genai
from google.generativeai import GenerativeModel
import requests


def text_extraction(ppt_file):
    presentation = Presentation(ppt_file)

    extracted_text = []

    for slides in presentation.slides:
        for shape in slides.shapes:
            if hasattr(shape, "text"):
                extracted_text.append(shape.text)

    combined_text = "\n".join(extracted_text)

    return combined_text

#apikey = "AIzaSyCJ88aztjeTldBwJCF4QxqrbyzLINqlzBY"
apikey="AIzaSyBmXC-dYqs0uVwonBYFYAuEPh0IzN-ACBE"

def generate_questions_geminai(text):
    # Replace "YOURAPIKEY" with your actual API key

    # Configure the SDK with your API key
    genai.configure(api_key=apikey)

    # Create a model instance for Gemini
    model = GenerativeModel("gemini-pro")

    # Generate text from a prompt
    prompt = "Write 5 multiple choice questions about " + text
    response = model.generate_content(prompt)
   # response = str(response)
    #data = response.json()
    print(response)
    print("-------------------------------------------")
    print(response.text)



def main():
    ppt_file = "slides/h.pptx"
    extracted_text = text_extraction(ppt_file)
    generate_questions_geminai(extracted_text)
    #for i, question in enumerate(questions, 1):
     #   print(f"Question {i}: {question}")

if __name__ == "__main__":
    main()
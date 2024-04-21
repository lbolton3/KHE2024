from pptx import Presentation
import google.generativeai as genai
from google.generativeai import GenerativeModel
import requests


def text_extraction(ppt_file):
    presentation = Presentation(ppt_file)

    extracted_text = []

    for slides in presentation.slides:
        for shape in slides.shapes:
            # print(slides.slide_id)
            if hasattr(shape, "text"):
                extracted_text.append(shape.text)

    combined_text = "\n".join(extracted_text)

    return combined_text


apikey = "AIzaSyBmXC-dYqs0uVwonBYFYAuEPh0IzN-ACBE" # yug's key

def generate_questions_geminai(text):
    # Replace "YOURAPIKEY" with your actual API key

    # Configure the SDK with your API key
    genai.configure(api_key=apikey)

    # Create a model instance for Gemini
    model = GenerativeModel("gemini-1.5-pro-latest")

    # Generate text from a prompt
    prompt = "Write 10 multiple choice questions with answers about the content of these slides and do not repeat questions included in slides and only create questions about content covered here - " + text
    response = model.generate_content(prompt)
   # response = str(response)
    #data = response.json()
    #print(response)
    print("-------------------------------------------")
    print(response.text)



def main():
    ppt_files = []

    # ppt_file = "./2cs1-ppt.pptx"
    # extracted_text = text_extraction(ppt_file)
    # print(extracted_text)
    # generate_questions_geminai(extracted_text)
    #for i, question in enumerate(questions, 1):
     #   print(f"Question {i}: {question}")
    
    # tempArr = ["1cs1-ppt.pptx"]
    for i in range(5):
        ppt_file = str(i) + "cs1-ppt.pptx"
        extracted_text = text_extraction(ppt_file)
        generate_questions_geminai(extracted_text)
        print("----------------------------------------------------------")

if __name__ == "__main__":
    main()